from collections.abc import Iterator
from importlib import metadata
from pathlib import Path
from typing import Any

from backend.ingestion.parser.errors import ParserDependencyError
from backend.ingestion.parser.models import ParsedBlock, ParsedDocument, ParsedPage


class PptxParser:
    """PPTX parser adapter backed by python-pptx."""

    parser_name = "python_pptx"
    supported_extensions = frozenset({".pptx"})
    supported_mime_types = frozenset(
        {"application/vnd.openxmlformats-officedocument.presentationml.presentation"}
    )

    def parse(self, file_path: Path) -> ParsedDocument:
        """Extract text boxes and table rows from every PPTX slide.

        Args:
            file_path: Absolute or relative path to a PPTX file stored by ingestion.

        Returns:
            A ParsedDocument with one ParsedPage per presentation slide.
        """
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise ParserDependencyError(self.parser_name, "python-pptx") from exc

        presentation = Presentation(file_path)
        pages: list[ParsedPage] = []

        # Create a page-like normalized record for each slide in presentation order.
        for slide_index, slide in enumerate(presentation.slides):
            page_number = slide_index + 1
            blocks = self._parse_slide_blocks(slide, page_number)
            slide_text = "\n\n".join(block.text for block in blocks)

            pages.append(
                ParsedPage(
                    page_number=page_number,
                    text=slide_text,
                    blocks=blocks,
                    metadata={"slide_number": page_number},
                )
            )

        return ParsedDocument(
            text="\n\n".join(page.text for page in pages).strip(),
            pages=pages,
            metadata={"source_path": str(file_path)},
            parser_name=self.parser_name,
            parser_version=self._get_parser_version(),
        )

    def _parse_slide_blocks(self, slide: Any, page_number: int) -> list[ParsedBlock]:
        """Convert supported text-bearing shapes on one slide into text blocks.

        Args:
            slide: Open python-pptx Slide instance to extract from.
            page_number: One-based slide number used as the normalized page number.

        Returns:
            A list of ParsedBlock instances in the slide shape order.
        """
        blocks: list[ParsedBlock] = []

        # Visit groups recursively so text inside grouped slide elements is retained.
        for shape_index, shape in enumerate(self._iter_shapes(slide.shapes)):
            if shape.has_text_frame:
                text = shape.text.strip()

                # Add non-empty native PowerPoint text as an independently traceable block.
                if text:
                    blocks.append(
                        ParsedBlock(
                            text=text,
                            page_number=page_number,
                            block_index=len(blocks),
                            metadata={
                                "block_type": "text",
                                "shape_index": shape_index,
                                "shape_id": shape.shape_id,
                            },
                        )
                    )

            # Add each non-empty table row separately to preserve table row context.
            if shape.has_table:
                self._append_table_rows(shape, page_number, shape_index, blocks)

        return blocks

    def _iter_shapes(self, shapes: Any) -> Iterator[Any]:
        """Yield regular shapes and the children of grouped shapes in display order.

        Args:
            shapes: python-pptx shape collection from a slide or group shape.

        Returns:
            An iterator over shapes that can contain text or tables.
        """
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        # Traverse every shape and recursively flatten group shapes.
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from self._iter_shapes(shape.shapes)
                continue

            yield shape

    def _append_table_rows(
        self,
        shape: Any,
        page_number: int,
        shape_index: int,
        blocks: list[ParsedBlock],
    ) -> None:
        """Append non-empty rows from one PPTX table shape to parsed blocks.

        Args:
            shape: python-pptx shape that owns the table being extracted.
            page_number: One-based slide number for the parsed blocks.
            shape_index: Position of the table shape among parsed slide shapes.
            blocks: Mutable parsed-block list that receives table row blocks.

        Returns:
            None. Rows are appended directly to the supplied block list.
        """
        # Emit one block per non-empty row so chunking can retain row provenance.
        for row_index, row in enumerate(shape.table.rows):
            cells = [cell.text.strip().replace("\n", " | ") for cell in row.cells]
            text = " | ".join(cell for cell in cells if cell)

            if text:
                blocks.append(
                    ParsedBlock(
                        text=text,
                        page_number=page_number,
                        block_index=len(blocks),
                        metadata={
                            "block_type": "table_row",
                            "shape_index": shape_index,
                            "shape_id": shape.shape_id,
                            "row_index": row_index,
                        },
                    )
                )

    def _get_parser_version(self) -> str:
        """Read the installed python-pptx package version.

        Args:
            None.

        Returns:
            The installed python-pptx version, or "unknown" when unavailable.
        """
        try:
            return metadata.version("python-pptx")
        except metadata.PackageNotFoundError:
            return "unknown"
