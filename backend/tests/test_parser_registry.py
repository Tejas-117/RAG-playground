import builtins
import tempfile
import unittest
from pathlib import Path
from unittest.mock import patch

from backend.ingestion.parsers import (
    DocxParser,
    MarkdownParser,
    ParserDependencyError,
    PlainTextParser,
    PptxParser,
    PyMuPDFParser,
    UnsupportedFileTypeError,
    build_default_parser_registry,
)


class ParserRegistryTestCase(unittest.TestCase):
    """Regression tests for ingestion parser selection and parser outputs."""

    def test_registry_selects_parser_by_extension_despite_mime_like_name(
        self,
    ) -> None:
        """Verify parser routing depends only on the filename extension.

        Returns:
            None. Assertions fail the test when parser selection changes.
        """
        registry = build_default_parser_registry()

        parser = registry.get_parser_by_extension(Path("application-pdf.txt"))

        self.assertIsInstance(parser, PlainTextParser)

    def test_registry_selects_parser_by_extension(
        self,
    ) -> None:
        """Verify extension routing selects a parser without upload metadata.

        Returns:
            None. Assertions fail the test when extension routing changes.
        """
        registry = build_default_parser_registry()

        parser = registry.get_parser_by_extension(Path("notes.md"))

        self.assertIsInstance(parser, MarkdownParser)

    def test_registry_selects_office_and_ebook_parsers_by_extension(self) -> None:
        """Verify the default registry selects every newly supported file type.

        Returns:
            None. Assertions fail the test when parser routing changes.
        """
        registry = build_default_parser_registry()

        # Check extension routing for the native Office parsers and PyMuPDF ebooks.
        self.assertIsInstance(
            registry.get_parser_by_extension(Path("report.docx")), DocxParser
        )
        self.assertIsInstance(
            registry.get_parser_by_extension(Path("slides.pptx")), PptxParser
        )
        self.assertIsInstance(
            registry.get_parser_by_extension(Path("book.epub")), PyMuPDFParser
        )
        self.assertIsInstance(
            registry.get_parser_by_extension(Path("book.mobi")), PyMuPDFParser
        )

    def test_registry_raises_for_unsupported_file_type(self) -> None:
        """Verify unsupported uploads fail with a structured parser error.

        Returns:
            None. Assertions fail the test when unsupported file handling changes.
        """
        registry = build_default_parser_registry()

        with self.assertRaises(UnsupportedFileTypeError):
            registry.get_parser_by_extension(Path("archive.zip"))

    def test_plain_text_parser_returns_normalized_document(self) -> None:
        """Verify plain text parsing returns the shared ParsedDocument model.

        Returns:
            None. Assertions fail the test when text parser output changes.
        """
        parser = PlainTextParser()

        with tempfile.TemporaryDirectory() as directory:
            file_path = Path(directory) / "notes.txt"
            file_path.write_text("hello\nworld", encoding="utf-8")

            parsed_document = parser.parse(file_path)

        self.assertEqual(parsed_document.text, "hello\nworld")
        self.assertEqual(parsed_document.parser_name, "plain_text")
        self.assertEqual(parsed_document.pages, [])
        self.assertEqual(parsed_document.metadata, {})

    def test_pymupdf_parser_raises_dependency_error_when_package_is_missing(
        self,
    ) -> None:
        """Verify PDF parser reports a missing optional dependency clearly.

        Returns:
            None. Assertions fail the test when PDF dependency errors change.
        """
        parser = PyMuPDFParser()
        original_import = builtins.__import__

        def import_without_fitz(name: str, *args: object, **kwargs: object) -> object:
            """Block fitz imports while delegating every other import.

            Args:
                name: Module name requested by Python import machinery.
                *args: Positional import arguments passed by Python.
                **kwargs: Keyword import arguments passed by Python.

            Returns:
                The imported module for non-fitz imports.
            """
            if name == "fitz":
                raise ImportError("No module named fitz")
            return original_import(name, *args, **kwargs)

        with (
            patch("builtins.__import__", side_effect=import_without_fitz),
            self.assertRaises(ParserDependencyError),
        ):
            parser.parse(Path("report.pdf"))

    def test_docx_parser_extracts_ordered_paragraph_and_table_blocks(self) -> None:
        """Verify DOCX parsing preserves paragraph and table-row content.

        Returns:
            None. Assertions fail the test when DOCX extraction changes.
        """
        from docx import Document

        parser = DocxParser()

        with tempfile.TemporaryDirectory() as directory:
            file_path = Path(directory) / "report.docx"
            document = Document()
            document.add_paragraph("Introduction")
            table = document.add_table(rows=1, cols=2)
            table.cell(0, 0).text = "Metric"
            table.cell(0, 1).text = "Score"
            document.add_paragraph("Conclusion")
            document.save(file_path)

            parsed_document = parser.parse(file_path)

        self.assertEqual(
            [block.text for block in parsed_document.pages[0].blocks],
            ["Introduction", "Metric | Score", "Conclusion"],
        )
        self.assertEqual(
            parsed_document.pages[0].blocks[1].metadata["block_type"],
            "table_row",
        )

    def test_pptx_parser_extracts_slide_text_and_table_rows(self) -> None:
        """Verify PPTX parsing maps slides, text boxes, and tables to parsed output.

        Returns:
            None. Assertions fail the test when PPTX extraction changes.
        """
        from pptx import Presentation
        from pptx.util import Inches

        parser = PptxParser()

        with tempfile.TemporaryDirectory() as directory:
            file_path = Path(directory) / "slides.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            text_box = slide.shapes.add_textbox(
                Inches(1), Inches(1), Inches(5), Inches(1)
            )
            text_box.text = "Experiment results"
            table = slide.shapes.add_table(
                1, 2, Inches(1), Inches(2), Inches(5), Inches(1)
            )
            table.table.cell(0, 0).text = "Recall"
            table.table.cell(0, 1).text = "0.90"
            presentation.save(file_path)

            parsed_document = parser.parse(file_path)

        self.assertEqual(len(parsed_document.pages), 1)
        self.assertEqual(
            [block.text for block in parsed_document.pages[0].blocks],
            ["Experiment results", "Recall | 0.90"],
        )
        self.assertEqual(
            parsed_document.pages[0].blocks[1].metadata["block_type"],
            "table_row",
        )


if __name__ == "__main__":
    unittest.main()
